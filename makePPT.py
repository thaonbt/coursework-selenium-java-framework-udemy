from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Project Presentation"
subtitle.text = "Objective, Test Scope, Example, Guidance, Demo"

# Slide 2: Objective
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Objective"
content.text = (
    "The objective of this project is to automate the testing of the registration and order submission processes "
    "using Selenium and Cucumber. The goal is to ensure that the application functions correctly and efficiently."
)

# Slide 3: Test Scope
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Test Scope"
content.text = (
    "The test scope includes:\n"
    "1. User registration process\n"
    "2. Order submission process\n"
    "3. Verification of order history\n"
    "4. Validation of error messages for invalid inputs"
)

# Slide 4: Example
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Example"
content.text = (
    "Example Test Case: Complete Registration with Valid Data\n\n"
    "Test Steps:\n"
    "1. Enter 'John' in the 'First Name' field.\n"
    "2. Enter 'Doe' in the 'Last Name' field.\n"
    "3. Enter 'john.doe@example.com' in the 'Email' field.\n"
    "4. Enter '1234567890' in the 'Phone Number' field.\n"
    "5. Select 'Engineer' from the 'Occupation' dropdown.\n"
    "6. Select 'Male' for the 'Gender' radio button.\n"
    "7. Enter 'password123' in the 'Password' field.\n"
    "8. Enter 'password123' in the 'Confirm Password' field.\n"
    "9. Check the 'I am 18 year or Older' checkbox.\n"
    "10. Click the 'Register' button.\n\n"
    "Expected Result: The user should see a success message indicating that the registration was successful."
)

# Slide 5: Guidance on How to Use
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Guidance on How to Use"
content.text = (
    "1. Clone the repository from GitHub.\n"
    "2. Install dependencies using Maven: `mvn test-compile`.\n"
    "3. Open the project in IntelliJ IDEA or your preferred IDE.\n"
    "4. Run the tests using Maven: `mvn test` or directly by TestNG.\n"
    "5. Review the test results and reports generated."
)

# Slide 6: Demo
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Demo"
content.text = (
    "The demo will showcase the following:\n"
    "1. User registration with valid data\n"
    "2. Order submission process\n"
    "3. Verification of order in order history\n"
    "4. Handling of invalid inputs and error messages"
)

# Save the presentation
prs.save('Project_Presentation.pptx')